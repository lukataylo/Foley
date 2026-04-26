"""Generate the Foley pitch deck (.pptx) — Apple-keynote economy on Foley's
own dark + aurora-amber palette. Eight slides, one sentence each, with a
single bento for the feature breakdown.

Slide flow:
  1. Title
  2. The problem (one sentence)
  3. What it is (one sentence)
  4. Live demo cue
  5. Feature bento (the only crowded slide)
  6. Built in 26h — what's new vs existing
  7. Who it's for / why now
  8. Close
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT_DIR = Path(__file__).parent
OUT_PPTX = OUT_DIR / "foley.pptx"

# ── Palette (matches apps/cutroom/src/app/globals.css dark theme) ──────
BG       = RGBColor(0x0A, 0x0A, 0x0A)
FG       = RGBColor(0xFA, 0xFA, 0xFA)
MUTED    = RGBColor(0x7A, 0x7A, 0x80)
AMBER    = RGBColor(0xF5, 0xB9, 0x4A)   # aurora-amber accent
AMBER_2  = RGBColor(0xC8, 0x88, 0x1F)
RED      = RGBColor(0xFF, 0x6B, 0x5C)
PANEL    = RGBColor(0x18, 0x18, 0x1A)
PANEL_2  = RGBColor(0x22, 0x22, 0x26)

# 16:9 widescreen.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Apple-style fonts. We use system fallbacks so the file renders even on
# machines without SF Pro Display installed.
FONT_DISPLAY = "SF Pro Display"
FONT_MONO    = "SF Mono"
FONT_FALLBACK = "Helvetica Neue"


def make_pres() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def fill_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_aurora_glow(slide, x, y, size, color=AMBER):
    """No-op. PowerPoint shapes are opaque; the layered-ovals trick to
    fake a CSS radial-gradient just paints concentric brown circles. The
    deck reads cleaner as pure-black-with-amber-labels Apple-keynote
    style; the amber section labels above each big sentence carry the
    brand presence on their own."""
    return


def add_text(
    slide, x, y, w, h, text,
    *, size=32, color=FG, bold=False, font=FONT_DISPLAY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.05,
    paragraph_gap=0,
):
    """Add a text box. If `text` contains \\n, each line becomes its own
    paragraph with tight spacing — LibreOffice was inserting a default
    paragraph gap that visibly broke the line rhythm of bullet lists."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(paragraph_gap)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def add_chrome(slide, page_num: int, total: int):
    """Top-right page counter + bottom-left brand mark, in the spirit of the
    hackathon's pagination header. Keeps the "monospace label" feel without
    being noisy."""
    add_text(
        slide,
        SLIDE_W - Inches(1.6), Inches(0.35), Inches(1.4), Inches(0.3),
        f"{page_num:02d} / {total:02d}",
        size=11, color=MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        Inches(0.6), Inches(0.35), Inches(3), Inches(0.3),
        "FOLEY · UNICORN MAFIA",
        size=11, color=MUTED, font=FONT_MONO,
    )


# ── Slide 1: Title ─────────────────────────────────────────────────────
def slide_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_aurora_glow(s, Inches(-1), Inches(-1), Inches(4))
    add_aurora_glow(s, Inches(10.5), Inches(5.5), Inches(3))
    add_chrome(s, 1, total)

    # Big "Foley." centered, with the subtitle parked safely below the
    # baseline. 140pt is tall enough to feel like a Keynote title without
    # overrunning the next text box.
    add_text(
        s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.2),
        "Foley.",
        size=140, color=FG, bold=True,
        line_spacing=1.0, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
        "Video as a build artifact.",
        size=28, color=MUTED, font=FONT_DISPLAY,
        line_spacing=1.0, align=PP_ALIGN.CENTER,
    )


# ── Slide 2: The problem ───────────────────────────────────────────────
def slide_problem(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_chrome(s, 2, total)

    add_text(
        s, Inches(1), Inches(0.95), Inches(8), Inches(0.5),
        "THE PROBLEM",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )
    # Smaller text + wider box so the line breaks naturally on its own
    # comma instead of getting chopped at "out of\ndate".
    add_text(
        s, Inches(1), Inches(2.5), Inches(11.3), Inches(4.5),
        "Every product walkthrough goes stale the moment you ship.",
        size=54, color=FG, line_spacing=1.15,
    )


# ── Slide 3: What it is ────────────────────────────────────────────────
def slide_what(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_aurora_glow(s, Inches(9), Inches(4.5), Inches(4))
    add_chrome(s, 3, total)

    add_text(
        s, Inches(1), Inches(0.95), Inches(8), Inches(0.5),
        "WHAT IT IS",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )
    add_text(
        s, Inches(1), Inches(2.5), Inches(11.3), Inches(4.5),
        "Foley rebuilds your walkthrough video from the codebase — every PR.",
        size=48, color=FG, line_spacing=1.2,
    )


# ── Slide 4: Live demo cue ─────────────────────────────────────────────
def slide_demo(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_aurora_glow(s, Inches(-2), Inches(2.5), Inches(8))
    add_chrome(s, 4, total)

    add_text(
        s, Inches(1), Inches(0.95), Inches(8), Inches(0.5),
        "LIVE",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )
    # Centered, with the arrow as a separate small accent line so the giant
    # display type doesn't have to squeeze a glyph that PowerPoint may not
    # have a font for at that size.
    add_text(
        s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(2.2),
        "Demo.",
        size=160, color=FG, bold=True,
        line_spacing=1.0, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
        "PR opens · video updates · we cut it in the editor · we publish.",
        size=20, color=MUTED, line_spacing=1.0, align=PP_ALIGN.CENTER,
    )


# ── Slide 5: Bento — feature breakdown ─────────────────────────────────
def slide_bento(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_chrome(s, 5, total)

    add_text(
        s, Inches(0.9), Inches(0.95), Inches(8), Inches(0.5),
        "WHAT'S IN THE BOX",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )

    # Bento grid: 3 cols x 2 rows. Each tile is a rounded panel with a glyph,
    # a title, and a one-line description.
    cells = [
        ("🎬", "Auto-capture", "Playwright records every step on every PR."),
        ("🎤", "AI narration",  "ElevenLabs reads the script in your brand voice."),
        ("🎵", "Music beds",    "Generated to fit length and tone — no licensing."),
        ("✨", "Title cards",   "Aurora gradient + screenshots + typed text."),
        ("🍌", "Nano Banana",   "Gemini composes mockups from real frames."),
        ("📤", "Multi-format",  "MP4, WebM, GIF, MP3, project bundle."),
    ]
    cols, rows = 3, 2
    margin_x = Inches(0.9)
    margin_top = Inches(1.7)
    gap = Inches(0.22)
    grid_w = SLIDE_W - margin_x * 2
    grid_h = SLIDE_H - margin_top - Inches(0.6)
    tile_w = (grid_w - gap * (cols - 1)) / cols
    tile_h = (grid_h - gap * (rows - 1)) / rows

    for i, (glyph, title, desc) in enumerate(cells):
        c, r = i % cols, i // cols
        x = margin_x + c * (tile_w + gap)
        y = margin_top + r * (tile_h + gap)
        tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tile_w, tile_h)
        tile.adjustments[0] = 0.06
        tile.fill.solid()
        tile.fill.fore_color.rgb = PANEL
        tile.line.color.rgb = PANEL_2
        tile.line.width = Pt(0.75)
        tile.shadow.inherit = False

        # Glyph
        add_text(
            s, x + Inches(0.5), y + Inches(0.35),
            tile_w - Inches(1), Inches(0.9),
            glyph, size=44, color=FG, line_spacing=1.0,
        )
        # Title
        add_text(
            s, x + Inches(0.5), y + Inches(1.4),
            tile_w - Inches(1), Inches(0.5),
            title, size=22, color=FG, bold=True, line_spacing=1.0,
        )
        # Description
        add_text(
            s, x + Inches(0.5), y + Inches(1.95),
            tile_w - Inches(1), Inches(1),
            desc, size=14, color=MUTED, line_spacing=1.25,
        )


# ── Slide 6: Built in 26h ──────────────────────────────────────────────
def slide_26h(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_chrome(s, 6, total)

    add_text(
        s, Inches(1), Inches(0.95), Inches(8), Inches(0.5),
        "BUILT IN 26H",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )

    # Two columns: existing vs new.
    col_y = Inches(2.2)
    col_h = Inches(4.5)

    # Existing column
    add_text(
        s, Inches(1), col_y, Inches(5.5), Inches(0.5),
        "EXISTING",
        size=11, color=MUTED, font=FONT_MONO, bold=True,
    )
    add_text(
        s, Inches(1), col_y + Inches(0.6), Inches(5.5), col_h,
        ("Director agent\n"
         "Playwright capture\n"
         "Walkthrough YAML\n"
         "Continuous narration synth\n"
         "Editor scaffolding"),
        size=22, color=MUTED, line_spacing=1.4, paragraph_gap=4,
    )

    # New column
    add_text(
        s, Inches(7), col_y, Inches(5.5), Inches(0.5),
        "NEW IN 26H",
        size=11, color=AMBER, font=FONT_MONO, bold=True,
    )
    add_text(
        s, Inches(7), col_y + Inches(0.6), Inches(5.5), col_h,
        ("Title cards composited into export\n"
         "Multi-format export (MP4 · WebM · GIF · MP3)\n"
         "Multi-select · split · duplicate · search\n"
         "Project bundle · YAML-aware undo\n"
         "20-fix polish + crash hardening"),
        size=22, color=FG, line_spacing=1.4, paragraph_gap=4,
    )


# ── Slide 7: Who / why now ─────────────────────────────────────────────
def slide_who(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_aurora_glow(s, Inches(9.5), Inches(-1.5), Inches(5))
    add_chrome(s, 7, total)

    add_text(
        s, Inches(1), Inches(0.95), Inches(8), Inches(0.5),
        "WHO · WHY NOW",
        size=12, color=AMBER, font=FONT_MONO, bold=True,
    )
    add_text(
        s, Inches(1), Inches(2.5), Inches(11.3), Inches(4.5),
        "For dev tools that ship daily — AI is finally good enough to narrate, mix, and composite without us.",
        size=42, color=FG, line_spacing=1.25,
    )


# ── Slide 8: Close ─────────────────────────────────────────────────────
def slide_close(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_aurora_glow(s, Inches(5.5), Inches(2.5), Inches(8))
    add_chrome(s, 8, total)

    add_text(
        s, Inches(1), Inches(2.6), Inches(11.3), Inches(2),
        "Your video stays current.",
        size=72, color=FG, bold=True, line_spacing=1.0,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(1), Inches(5), Inches(11.3), Inches(0.6),
        "github.com/lukataylo/Foley",
        size=22, color=AMBER, font=FONT_MONO,
        align=PP_ALIGN.CENTER, line_spacing=1.0,
    )


def main() -> None:
    prs = make_pres()
    total = 8
    slide_title(prs, total)
    slide_problem(prs, total)
    slide_what(prs, total)
    slide_demo(prs, total)
    slide_bento(prs, total)
    slide_26h(prs, total)
    slide_who(prs, total)
    slide_close(prs, total)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
